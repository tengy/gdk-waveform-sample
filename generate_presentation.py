"""
generate_presentation.py
Generates a 10-slide PowerPoint presentation covering "Introduction to Multimedia"
based on the Multimedia.txt textbook content.

Requirements:
    pip install python-pptx

Usage:
    python generate_presentation.py
    --> Outputs: Introduction_to_Multimedia.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ────────────────────────────────────────────────────────────
C_DARK_BLUE  = RGBColor(0x1A, 0x37, 0x6C)
C_MID_BLUE   = RGBColor(0x2E, 0x6D, 0xB4)
C_LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF7)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
C_ORANGE     = RGBColor(0xF0, 0x8C, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    fill = shape.fill
    if fill_rgb:
        fill.solid()
        fill.fore_color.rgb = fill_rgb
    else:
        fill.background()
    line = shape.line
    if line_rgb:
        line.color.rgb = line_rgb
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text, font_size=18, bold=False,
                color=C_DARK_TEXT, align=PP_ALIGN.LEFT, italic=False,
                word_wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_slide_base(title_text, subtitle_text="", chapter_tag=""):
    """Adds a new blank slide with a dark-blue header bar and title."""
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF4, 0xF7, 0xFB))
    add_rect(slide, 0, 0, 13.33, 1.35, fill_rgb=C_DARK_BLUE)
    if chapter_tag:
        add_textbox(slide, 10.2, 0.08, 3.0, 0.35, chapter_tag,
                    font_size=9, color=C_LIGHT_BLUE, align=PP_ALIGN.RIGHT)
    add_textbox(slide, 0.35, 0.12, 12.6, 0.85, title_text,
                font_size=28, bold=True, color=C_WHITE)
    if subtitle_text:
        add_textbox(slide, 0.35, 1.4, 12.6, 0.45, subtitle_text,
                    font_size=13, italic=True, color=C_MID_BLUE)
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_rgb=C_MID_BLUE)
    add_textbox(slide, 0.2, 7.12, 12.9, 0.3,
                "Introduction to Multimedia  -  Dr. Ahmed H. Eltengy",
                font_size=8, color=C_WHITE, align=PP_ALIGN.CENTER)
    return slide


def add_bullet_box(slide, l, t, w, h, bullets, font_size=14,
                   color=C_DARK_TEXT, title=None, title_color=C_MID_BLUE,
                   bg_color=None, indent_char="- "):  
    """Draws an optional titled bullet list inside a coloured card."""
    if bg_color:
        add_rect(slide, l, t, w, h, fill_rgb=bg_color,
                 line_rgb=C_MID_BLUE, line_width_pt=0.75)
    top = t + 0.05
    if title:
        add_textbox(slide, l + 0.15, top, w - 0.3, 0.38, title,
                    font_size=13, bold=True, color=title_color)
        top += 0.38
    txb = slide.shapes.add_textbox(
        Inches(l + 0.15), Inches(top),
        Inches(w - 0.3), Inches(h - (top - t) - 0.1)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = indent_char + b
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def add_simple_table(slide, l, t, w, headers, rows,
                     header_font=11, body_font=10, col_widths=None):
    """Draws a styled table from plain Python lists."""
    if col_widths is None:
        col_widths = [w / len(headers)] * len(headers)
    ROW_H = 0.32
    hdr_h = 0.38
    x = l
    for hdr, cw in zip(headers, col_widths):
        add_rect(slide, x, t, cw, hdr_h, fill_rgb=C_MID_BLUE,
                 line_rgb=C_WHITE, line_width_pt=0.5)
        add_textbox(slide, x + 0.05, t + 0.03, cw - 0.1, hdr_h - 0.06,
                    hdr, font_size=header_font, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        x += cw
    for ri, row in enumerate(rows):
        y = t + hdr_h + ri * ROW_H
        bg = C_LIGHT_BLUE if ri % 2 == 0 else C_WHITE
        x = l
        for cell, cw in zip(row, col_widths):
            add_rect(slide, x, y, cw, ROW_H, fill_rgb=bg,
                     line_rgb=RGBColor(0xCC, 0xD9, 0xEA), line_width_pt=0.4)
            add_textbox(slide, x + 0.06, y + 0.03, cw - 0.12, ROW_H - 0.06,
                        str(cell), font_size=body_font, color=C_DARK_TEXT)


def add_key_box(slide, l, t, w, h, label, text, font_size=11):
    """Orange-left-border key definition box."""
    add_rect(slide, l, t, 0.07, h, fill_rgb=C_ORANGE)
    add_rect(slide, l + 0.07, t, w - 0.07, h,
             fill_rgb=RGBColor(0xFF, 0xF3, 0xD0),
             line_rgb=C_ORANGE, line_width_pt=0.6)
    add_textbox(slide, l + 0.18, t + 0.05, w - 0.28, 0.3,
                label, font_size=11, bold=True, color=C_ORANGE)
    add_textbox(slide, l + 0.18, t + 0.33, w - 0.28, h - 0.38,
                text, font_size=font_size, color=C_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 -- Title / Overview
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=C_DARK_BLUE)

for cx, cy, r, col in [
    (11.5, 1.0, 1.8, RGBColor(0x2E, 0x6D, 0xB4)),
    (1.2,  6.2, 1.2, RGBColor(0x1E, 0x4D, 0x94)),
    (12.5, 6.5, 0.9, RGBColor(0x1E, 0x4D, 0x94)),
]:
    shape = slide.shapes.add_shape(9, Inches(cx - r), Inches(cy - r),
                                   Inches(r * 2), Inches(r * 2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = col
    shape.line.fill.background()

add_textbox(slide, 1.5, 1.5, 10.3, 1.1, "Introduction to Multimedia",
            font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 3.5, 2.75, 6.33, 0.06, fill_rgb=C_ORANGE)
add_textbox(slide, 1.5, 2.9, 10.3, 0.55,
            "A Complete University Textbook  -  Undergraduate / 1st Year",
            font_size=16, color=C_LIGHT_BLUE, align=PP_ALIGN.CENTER)

chapters = ["Foundations", "Graphics", "Color & Video",
            "Animation", "Audio", "Typography", "Design", "Accessibility"]
pill_w, pill_h, pill_gap = 1.4, 0.36, 0.12
total_w = len(chapters) * pill_w + (len(chapters) - 1) * pill_gap
start_x = (13.33 - total_w) / 2
for i, ch in enumerate(chapters):
    px = start_x + i * (pill_w + pill_gap)
    add_rect(slide, px, 3.65, pill_w, pill_h,
             fill_rgb=C_MID_BLUE, line_rgb=C_LIGHT_BLUE, line_width_pt=0.5)
    add_textbox(slide, px, 3.65, pill_w, pill_h,
                ch, font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.3, 10.3, 0.45, "Dr. Ahmed H. Eltengy",
            font_size=18, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
add_rect(slide, 0, 7.1, 13.33, 0.4, fill_rgb=RGBColor(0x0F, 0x22, 0x44))
add_textbox(slide, 0.2, 7.12, 12.9, 0.3,
            "10-Chapter Comprehensive Course  |  Covers Chapters 1-10",
            font_size=9, color=RGBColor(0xAA, 0xC4, 0xE8), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 -- Chapter 1
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide_base(
    "Chapter 1 - Introduction to Multimedia",
    "What is Multimedia? - History - Components - Production Pipeline - Digitization",
    "Chapter 1"
)

add_key_box(slide, 0.3, 1.85, 5.8, 1.0,
            "KEY DEFINITION",
            "Multimedia: Integration of 2+ media types (text, images, audio, "
            "video, animation) presented through a digital, interactive system.")

add_bullet_box(slide, 0.3, 3.0, 5.8, 3.0,
               ["Text -- character codes (Unicode)",
                "Image -- pixel grids (bitmap)",
                "Audio -- sampled amplitude values",
                "Video -- sequence of image frames",
                "Animation -- keyframes + interpolation",
                "3D Models -- vertices, edges, faces"],
               title="6 Core Media Types",
               bg_color=RGBColor(0xEB, 0xF2, 0xFB), font_size=12)

add_textbox(slide, 6.3, 1.85, 6.8, 0.35,
            "Multimedia History Timeline",
            font_size=13, bold=True, color=C_MID_BLUE)
add_simple_table(
    slide, 6.3, 2.25, 6.8,
    headers=["Era", "Key Development", "Impact"],
    rows=[
        ["1980s", "CD-ROM, HyperCard", "Encyclopedias on disc"],
        ["1990s", "WWW, MPEG, Flash", "Global multimedia via internet"],
        ["2000s", "YouTube, iTunes", "User-generated content"],
        ["2010s", "HTML5, Mobile, 4K", "Plugin-free web, smartphones"],
        ["2020s", "WebXR, AI media, VR/AR", "Immersive & AI-assisted media"],
    ],
    col_widths=[0.85, 2.45, 3.5], header_font=10, body_font=9
)
add_bullet_box(slide, 6.3, 4.75, 6.8, 1.35,
               ["Plan -> Design -> Produce -> Post-produce -> Deliver",
                "Digitization: Sampling -> Quantization -> Encoding",
                "Nyquist theorem: sample rate >= 2x highest frequency"],
               title="Production Pipeline & Digitization",
               bg_color=RGBColor(0xEB, 0xF2, 0xFB), font_size=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 -- Chapter 2
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide_base(
    "Chapter 2 - Digital Media & Vector Graphics",
    "Pixel Grids - Vector vs Bitmap - Objects & Paths - Transformations",
    "Chapter 2"
)

for lx, title_txt, items, bg in [
    (0.3, "VECTOR GRAPHICS",
     ["Stored as mathematical geometry",
      "Resolution-independent (scales perfectly)",
      "Best for: logos, icons, diagrams",
      "Formats: SVG, AI, EPS, PDF",
      "Editable: fill, stroke, gradient",
      "Transforms: translate, scale, rotate, shear"],
     RGBColor(0xE8, 0xF4, 0xE8)),
    (6.8, "BITMAP / RASTER",
     ["Stored as fixed grid of pixels",
      "Resolution-dependent (pixelates on zoom)",
      "Best for: photos, textures",
      "Formats: JPEG, PNG, TIFF, GIF",
      "Described by Width x Height x Bit-depth",
      "Common sizes: HD 1280x720, 4K 3840x2160"],
     RGBColor(0xF4, 0xE8, 0xE8)),
]:
    add_bullet_box(slide, lx, 1.85, 6.2, 3.0,
                   items, title=title_txt,
                   bg_color=bg, font_size=12, title_color=C_DARK_BLUE)

add_textbox(slide, 0.3, 5.05, 12.7, 0.32,
            "Geometric Transformations (Chapter 2.7)",
            font_size=13, bold=True, color=C_MID_BLUE)

tforms = [
    ("Translation", "Moves object by (tx, ty);\nno size/orientation change"),
    ("Scaling",     "Resizes by sx, sy;\nunequal -> distortion"),
    ("Rotation",    "Turns by angle theta\naround a center point"),
    ("Shearing",    "Slants along one axis;\nother axis fixed"),
    ("Reflection",  "Mirror image across\nan axis (x, y, y=x)"),
]
card_w = 2.45
for i, (name, desc) in enumerate(tforms):
    cx = 0.3 + i * (card_w + 0.07)
    add_rect(slide, cx, 5.42, card_w, 1.55,
             fill_rgb=C_LIGHT_BLUE, line_rgb=C_MID_BLUE, line_width_pt=0.5)
    add_textbox(slide, cx + 0.1, 5.46, card_w - 0.2, 0.32,
                name, font_size=11, bold=True, color=C_DARK_BLUE)
    add_textbox(slide, cx + 0.1, 5.8, card_w - 0.2, 1.1,
                desc, font_size=9, color=C_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 -- Chapter 3
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide_base(
    "Chapter 3 - Graphics & Image Processing",
    "3D Pipeline - Raster Structure - Resolution - Compression - File Formats",
    "Chapter 3"
)

add_bullet_box(slide, 0.3, 1.85, 4.0, 2.0,
               ["Model Space -> World Space -> Camera Space",
                "Projection -> Screen Space -> Rasterize",
                "Shading: Flat -> Gouraud -> Phong -> PBR",
                "Hidden Surface: Z-buffer algorithm",
                "Texturing via UV coordinate mapping"],
               title="3D Rendering Pipeline",
               bg_color=RGBColor(0xEB, 0xF2, 0xFB), font_size=11)

add_textbox(slide, 4.5, 1.85, 8.6, 0.3,
            "Color Depth (Bit Depth)", font_size=12, bold=True, color=C_MID_BLUE)
add_simple_table(
    slide, 4.5, 2.2, 8.6,
    headers=["Bit Depth", "Colors", "Bytes/Pixel", "Use Case"],
    rows=[
        ["8-bit",  "256 (indexed)",    "1 B", "GIF, icons"],
        ["24-bit", "16.7 M (8+8+8)",   "3 B", "Standard web & print photos"],
        ["32-bit", "16.7 M + alpha",   "4 B", "PNG with transparency"],
        ["48-bit", "Trillions (16x3)", "6 B", "RAW professional photography"],
    ],
    col_widths=[1.6, 2.4, 1.5, 3.1], header_font=10, body_font=9
)
add_bullet_box(slide, 0.3, 4.0, 6.0, 2.1,
               ["Spatial redundancy: neighbouring pixels similar",
                "Perceptual redundancy: eye less sensitive to colour detail",
                "LOSSLESS (PNG, FLAC): perfect reconstruction; 2:1-5:1",
                "LOSSY (JPEG, H.264): data discarded; 10:1-100:1+",
                "RLE -> replace repeating runs with count+value",
                "LZW/Huffman: dictionary & frequency-based coding"],
               title="Why Compression? - Key Techniques",
               bg_color=RGBColor(0xFFF3D0), font_size=11)
add_simple_table(
    slide, 6.4, 4.0, 6.7,
    headers=["Format", "Type", "Best Use"],
    rows=[
        ["JPEG", "Lossy raster",    "Photographs, web/print"],
        ["PNG",  "Lossless raster", "Transparency, icons, UI"],
        ["SVG",  "Vector",          "Logos, scalable web graphics"],
        ["WebP", "Lossy/Lossless",  "Modern web - beats JPEG+PNG"],
        ["TIFF", "Lossless raster", "Professional print, archival"],
        ["RAW",  "Minimal loss",    "Camera originals, max flex."],
    ],
    col_widths=[1.1, 1.9, 3.7], header_font=10, body_font=9
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 -- Chapter 4
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide_base(
    "Chapter 4 - Color Science & Video",
    "Color Models - Color Management - Video Standards - Codecs - Editing",
    "Chapter 4"
)

add_simple_table(
    slide, 0.3, 1.85, 6.8,
    headers=["Model", "Type", "Primary Use"],
    rows=[
        ["RGB",     "Additive",     "Screens, web, digital cameras"],
        ["CMYK",    "Subtractive",  "Commercial printing"],
        ["HSV/HSL", "Cylindrical",  "Intuitive artistic color picking"],
        ["CIE Lab", "Perceptual",   "Device-independent color mgmt"],
        ["sRGB",    "Standardised", "Web standard, consumer devices"],
        ["DCI-P3",  "Wide gamut",   "Digital cinema, Apple displays"],
    ],
    col_widths=[1.4, 1.8, 3.6], header_font=10, body_font=9
)
add_textbox(slide, 7.3, 1.85, 5.8, 0.3,
            "Major Video Codecs", font_size=12, bold=True, color=C_MID_BLUE)
add_simple_table(
    slide, 7.3, 2.2, 5.8,
    headers=["Codec", "Year", "Key Application"],
    rows=[
        ["H.264 / AVC",  "2003", "Current web & Blu-ray standard"],
        ["H.265 / HEVC", "2013", "4K streaming; 40-50% better"],
        ["VP9",          "2013", "YouTube; royalty-free"],
        ["AV1",          "2018", "YouTube 4K; best compression"],
        ["Apple ProRes", "2007", "Professional editing (intraframe)"],
    ],
    col_widths=[2.0, 0.8, 3.0], header_font=10, body_font=9
)
add_bullet_box(slide, 0.3, 4.7, 6.8, 1.45,
               ["Y = Luminance  |  Cb/Cr = Colour difference",
                "4:4:4 -- Full resolution, max quality (RAW/professional)",
                "4:2:2 -- Cb/Cr at 1/2 horizontal -> 33% saving (broadcast)",
                "4:2:0 -- Cb/Cr at 1/2 H+V -> 50% saving (H.264, streaming)"],
               title="YCbCr & Chroma Subsampling",
               bg_color=RGBColor(0xEB, 0xF2, 0xFB), font_size=11)
add_bullet_box(slide, 7.3, 4.7, 5.8, 1.45,
               ["NLE: timeline-based, non-destructive editing",
                "Cut / Dissolve / Wipe / L-cut / J-cut / Match cut",
                "Compositing: alpha channel + chroma keying",
                "Green screen: green furthest from skin tone",
                "Color grading: LUTs, curves, lift/gamma/gain"],
               title="Editing & Compositing",
               bg_color=RGBColor(0xEB, 0xF2, 0xFB), font_size=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 -- Chapter 5
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide_base(
    "Chapter 5 - Animation Fundamentals",
    "12 Principles - Keyframes - Interpolation - Motion Graphics - Delivery",
    "Chapter 5"
)

principles = [
    ("1. Squash & Stretch",  "Deform to show mass/elasticity"),
    ("2. Anticipation",      "Small prep before main action"),
    ("3. Staging",           "Clear silhouette & composition"),
    ("4. S-Ahead/Pose-Pose", "Frame-by-frame vs key poses"),
    ("5. Follow-Through",    "Secondary parts continue moving"),
    ("6. Slow In/Slow Out",  "Natural acceleration/deceleration"),
    ("7. Arcs",              "Natural motion follows curves"),
    ("8. Secondary Action",  "Supporting details add richness"),
    ("9. Timing",            "Frames count = speed & weight"),
    ("10. Exaggeration",     "Amplify beyond realism for impact"),
    ("11. Solid Drawing",    "Form, weight, 3D feel in 2D"),
    ("12. Appeal",           "Watchable, personality, clear design"),
]
COLS, cell_w, cell_h = 4, 3.1, 0.84
for i, (name, desc) in enumerate(principles):
    row, col = divmod(i, COLS)
    cx = 0.3 + col * (cell_w + 0.07)
    cy = 1.85 + row * (cell_h + 0.06)
    add_rect(slide, cx, cy, cell_w, cell_h,
             fill_rgb=C_LIGHT_BLUE if (row + col) % 2 == 0 else C_WHITE,
             line_rgb=C_MID_BLUE, line_width_pt=0.4)
    add_textbox(slide, cx + 0.1, cy + 0.04, cell_w - 0.2, 0.3,
                name, font_size=9, bold=True, color=C_DARK_BLUE)
    add_textbox(slide, cx + 0.1, cy + 0.35, cell_w - 0.2, 0.44,
                desc, font_size=8, color=C_DARK_TEXT)

add_textbox(slide, 0.3, 5.7, 12.7, 0.3,
            "Animation Delivery Formats", font_size=12, bold=True, color=C_MID_BLUE)
fmt_data = [
    ("GIF",       "256 colors\n1-bit alpha\nSimple web anim."),
    ("APNG",      "24-bit\nFull alpha\nBetter GIF alt."),
    ("WebP",      "Full + alpha\nSmaller than GIF\nModern browsers"),
    ("Lottie",    "Vector JSON\nAfter Effects->web\nFull alpha"),
    ("CSS Anim.", "Code-based\nHW-accelerated\nUI micro-anim."),
    ("MP4/WebM",  "Full video\nComplex anim.\nLoop w/ player"),
]
fw = 2.1
for i, (name, desc) in enumerate(fmt_data):
    fx = 0.3 + i * (fw + 0.07)
    add_rect(slide, fx, 6.05, fw, 0.98,
             fill_rgb=RGBColor(0xE0, 0xED, 0xF9),
             line_rgb=C_MID_BLUE, line_width_pt=0.4)
    add_textbox(slide, fx + 0.08, 6.08, fw - 0.16, 0.3,
                name, font_size=10, bold=True, color=C_DARK_BLUE)
    add_textbox(slide, fx + 0.08, 6.4, fw - 0.16, 0.6,
                desc, font_size=8, color=C_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 -- Chapter 6
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide_base(
    "Chapter 6 - Audio",
    "Sound Properties - ADC Process - Signal Processing - Compression - MIDI",
    "Chapter 6"
)

add_bullet_box(slide, 0.3, 1.85, 4.0, 2.55,
               ["Frequency (Hz) -> perceived as Pitch",
                "Human range: ~20 Hz - 20,000 Hz",
                "Amplitude -> perceived as Loudness (dB)",
                "10 dB = x10 intensity, ~x2 perceived loud",
                "Timbre: determined by overtone structure",
                "ADC: Sampling -> Quantization -> Encoding",
                "Nyquist: sample rate >= 2x highest frequency"],
               title="Sound Properties & ADC",
               bg_color=RGBColor(0xEB, 0xF2, 0xFB), font_size=11)

add_simple_table(
    slide, 4.5, 1.85, 8.6,
    headers=["Sample Rate", "Nyquist Freq.", "Application"],
    rows=[
        ["8,000 Hz",   "4 kHz",   "Telephone / voice codecs"],
        ["44,100 Hz",  "22 kHz",  "CD audio -- standard consumer music"],
        ["48,000 Hz",  "24 kHz",  "Pro audio & video production"],
        ["96,000 Hz",  "48 kHz",  "High-resolution studio recording"],
        ["192,000 Hz", "96 kHz",  "Mastering / archiving"],
    ],
    col_widths=[2.1, 1.8, 4.7], header_font=10, body_font=9
)
add_bullet_box(slide, 0.3, 4.55, 6.2, 1.55,
               ["EQ: high-pass, low-pass, shelf, peak/bell, parametric",
                "Compressor: reduces gain above threshold",
                "Limiter: hard compression; prevents clipping",
                "Gate: reduces gain below threshold (removes bleed)",
                "Reverb: simulates space reflections (room, hall, plate)",
                "Delay: discrete echoes; chorus; flanger; phaser"],
               title="Signal Processing & Effects",
               bg_color=RGBColor(0xFFF3D0), font_size=11)
add_bullet_box(slide, 6.6, 4.55, 6.5, 1.55,
               ["MIDI: performance instructions, NOT audio samples",
                "Standardised 1983; still universal in music production",
                "Note On/Off, Velocity, Control Change, Pitch Bend",
                "Advantages: tiny files, fully editable, pitch/tempo free",
                "Formats: MP3, AAC, FLAC, WAV, Opus (best low-bitrate)"],
               title="MIDI & Audio Formats",
               bg_color=RGBColor(0xFFF3D0), font_size=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 -- Chapter 7
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide_base(
    "Chapter 7 - Text, Typography & The Web",
    "Character Encoding - Font Technologies - Typographic Fundamentals - Hypermedia",
    "Chapter 7"
)

add_simple_table(
    slide, 0.3, 1.85, 5.5,
    headers=["Encoding", "Bytes/Char", "Best For"],
    rows=[
        ["ASCII",  "1 (7-bit)", "English only; 128 characters"],
        ["UTF-8",  "1-4",       "Web standard; ASCII-compatible"],
        ["UTF-16", "2 or 4",    "Windows APIs, Java, JavaScript"],
        ["UTF-32", "4 (fixed)", "Simplicity; less common"],
    ],
    col_widths=[1.2, 1.3, 3.0], header_font=10, body_font=9
)
add_key_box(slide, 0.3, 3.5, 5.5, 0.85,
            "UTF-8 DOMINANCE",
            "Over 98% of web pages use UTF-8. "
            "Always declare: <meta charset=\"UTF-8\">")

add_simple_table(
    slide, 6.0, 1.85, 7.1,
    headers=["Technology", "Format", "Key Feature"],
    rows=[
        ["TrueType",       ".ttf",      "Single file; widely supported since 1991"],
        ["OpenType",       ".otf/.ttf", "Unicode; advanced features; supersedes TT"],
        ["WOFF2",          ".woff2",    "Compressed web fonts; ~30% smaller"],
        ["Variable Fonts", ".ttf/.otf", "One file, multiple axes (weight/width/italic)"],
    ],
    col_widths=[2.0, 1.4, 3.7], header_font=10, body_font=9
)
add_bullet_box(slide, 6.0, 3.5, 7.1, 2.6,
               ["Old Style Serif -- Garamond, Caslon -- books, academic",
                "Transitional Serif -- Times New Roman -- newspapers",
                "Geometric Sans -- Futura, Bauhaus -- modernist / tech",
                "Neo-Grotesque -- Helvetica, Arial -- signage, UI",
                "Monospace -- Consolas, Courier -- code, tables",
                "Script -- Pacifico, Brush Script -- invitations, logos"],
               title="Typeface Classifications",
               bg_color=RGBColor(0xEB, 0xF2, 0xFB), font_size=11)
add_bullet_box(slide, 0.3, 4.45, 5.5, 1.65,
               ["Hypertext: non-linear linked text (WWW foundation)",
                "Hypermedia: hypertext + images, audio, video",
                "Navigation: clear affordances, consistent structure",
                "Screen vs Print: lower resolution, light-emitting, scrolling",
                "Accessible: sufficient contrast, scalable text, focus order"],
               title="Hypermedia & Web Typography",
               bg_color=RGBColor(0xFFF3D0), font_size=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 -- Chapters 8 & 9
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide_base(
    "Chapters 8 & 9 - Visual Design & Interactivity",
    "Gestalt Psychology - Color in Design - Grid Systems - Interactive Storytelling",
    "Ch. 8-9"
)

add_bullet_box(slide, 0.3, 1.85, 6.1, 2.55,
               ["Proximity -- nearby elements perceived as related",
                "Similarity -- similar elements grouped together",
                "Continuity -- eye follows lines and curves",
                "Closure -- brain completes incomplete shapes",
                "Figure/Ground -- separation of object from background",
                "Common Fate -- elements moving together are grouped"],
               title="Gestalt Psychology Principles",
               bg_color=RGBColor(0xEB, 0xF2, 0xFB), font_size=12)
add_bullet_box(slide, 0.3, 4.5, 6.1, 1.65,
               ["Hue, Saturation, Value drive emotional response",
                "Warm colours (red/orange): energy, urgency, appetite",
                "Cool colours (blue/green): calm, trust, professionalism",
                "60-30-10 rule: dominant / secondary / accent colours",
                "Contrast ratio >= 4.5:1 for readable text (WCAG AA)"],
               title="Color in Visual Communication",
               bg_color=RGBColor(0xFFF3D0), font_size=11)
add_bullet_box(slide, 6.5, 1.85, 6.55, 2.55,
               ["Grid: invisible framework of columns & gutters",
                "Provides visual order, alignment, consistency",
                "Column grids (12-col most common in web design)",
                "Baseline grid: vertical rhythm for typography",
                "Rule of thirds: divide canvas 3x3, place focal points",
                "White space: breathing room, not 'wasted' space"],
               title="Grid Systems in Design",
               bg_color=RGBColor(0xEB, 0xF2, 0xFB), font_size=12)
add_bullet_box(slide, 6.5, 4.5, 6.55, 1.65,
               ["Interactivity: user controls flow and content",
                "Branching narratives -- choices alter story path",
                "Levels of interactivity: reactive -> contributory",
                "UX principles: feedback, affordances, discoverability",
                "Interactive storytelling: games, VR, choose-your-own"],
               title="Interactivity & Storytelling",
               bg_color=RGBColor(0xFFF3D0), font_size=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 -- Chapter 10
# ════════════════════════════════════════════════════════════���═════════════════
slide = add_slide_base(
    "Chapter 10 - Accessibility & Key Takeaways",
    "WCAG Guidelines - Alt Text - Captions - Inclusive Design - Emerging Tech",
    "Chapter 10"
)

add_simple_table(
    slide, 0.3, 1.85, 6.0,
    headers=["WCAG Principle", "What It Means"],
    rows=[
        ["Perceivable",    "Content available to all senses; alt text, captions"],
        ["Operable",       "All functions usable via keyboard; no seizure risk"],
        ["Understandable", "Clear language; predictable navigation; error help"],
        ["Robust",         "Compatible with assistive technologies (ARIA roles)"],
    ],
    col_widths=[2.2, 3.8], header_font=10, body_font=10
)
add_bullet_box(slide, 0.3, 3.85, 6.0, 2.3,
               ["Text alternatives: alt attribute for all meaningful images",
                "Video: closed captions + audio description track",
                "Contrast ratio >= 4.5:1 (text) / 3:1 (large text)",
                "Keyboard navigable: logical tab order + focus styles",
                "ARIA landmarks & roles for screen-reader structure",
                "Avoid relying on colour alone to convey meaning"],
               title="Practical Accessibility Checklist",
               bg_color=RGBColor(0xE8, 0xF4, 0xE8), font_size=11)
add_bullet_box(slide, 6.5, 1.85, 6.55, 1.65,
               ["WebXR & VR/AR: spatial, haptic, voice-driven UI",
                "AI-generated media: deepfakes, synthetic voices -> ethics",
                "Real-time captions: AI speech-to-text in live events",
                "Personalisation: adaptive interfaces (font size, contrast)"],
               title="Emerging Technologies",
               bg_color=RGBColor(0xEB, 0xF2, 0xFB), font_size=11)

summary = [
    ("Ch.1",  "Multimedia Foundations\nBinary, ADC, pipelines"),
    ("Ch.2",  "Digital Media & Vectors\nPixels, transformations"),
    ("Ch.3",  "Graphics & Compression\n3D pipeline, file formats"),
    ("Ch.4",  "Color Science & Video\nModels, codecs, NLE"),
    ("Ch.5",  "Animation\n12 principles, keyframes"),
    ("Ch.6",  "Audio\nSound, ADC, MIDI, EQ"),
    ("Ch.7",  "Typography & Web\nUnicode, fonts, hypertext"),
    ("Ch.8-9","Design & Interactivity\nGestalt, grids, UX"),
    ("Ch.10", "Accessibility\nWCAG, inclusive design"),
]
add_textbox(slide, 6.5, 3.6, 6.55, 0.3,
            "Course at a Glance", font_size=12, bold=True, color=C_MID_BLUE)
w2, h2 = 2.1, 0.75
for i, (ch, desc) in enumerate(summary):
    row, col = divmod(i, 3)
    cx2 = 6.5 + col * (w2 + 0.07)
    cy2 = 3.95 + row * (h2 + 0.06)
    add_rect(slide, cx2, cy2, w2, h2,
             fill_rgb=C_LIGHT_BLUE if (row + col) % 2 == 0 else C_WHITE,
             line_rgb=C_MID_BLUE, line_width_pt=0.4)
    add_textbox(slide, cx2 + 0.1, cy2 + 0.04, w2 - 0.2, 0.3,
                ch, font_size=9, bold=True, color=C_DARK_BLUE)
    add_textbox(slide, cx2 + 0.1, cy2 + 0.36, w2 - 0.2, 0.35,
                desc, font_size=8, color=C_DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
OUTPUT = "Introduction_to_Multimedia.pptx"
prs.save(OUTPUT)
print(f"Presentation saved -> {OUTPUT}")
print(f"   Slides: {len(prs.slides)}")